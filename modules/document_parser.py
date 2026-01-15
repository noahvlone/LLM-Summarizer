"""
Document parsing module.
Handles extraction of text from PDF and PowerPoint files.
"""
import io
import fitz  # PyMuPDF
from pptx import Presentation


def parse_pdf(file) -> dict:
    """
    Extract text content from a PDF file.
    
    Args:
        file: File object (from Streamlit uploader) or file path
        
    Returns:
        Dictionary with 'success', 'text', 'pages', and optional 'error' keys
    """
    try:
        # Handle Streamlit UploadedFile
        if hasattr(file, 'read'):
            file_bytes = file.read()
            doc = fitz.open(stream=file_bytes, filetype="pdf")
        else:
            doc = fitz.open(file)
        
        text_content = []
        
        for page_num, page in enumerate(doc, 1):
            page_text = page.get_text()
            if page_text.strip():
                text_content.append(f"[Page {page_num}]\n{page_text}")
        
        doc.close()
        
        full_text = "\n\n".join(text_content)
        
        return {
            'success': True,
            'text': full_text,
            'pages': len(text_content)
        }
        
    except Exception as e:
        return {
            'success': False,
            'text': '',
            'pages': 0,
            'error': f'Error parsing PDF: {str(e)}'
        }


def parse_pptx(file) -> dict:
    """
    Extract text content from a PowerPoint file.
    
    Args:
        file: File object (from Streamlit uploader) or file path
        
    Returns:
        Dictionary with 'success', 'text', 'slides', and optional 'error' keys
    """
    try:
        # Handle Streamlit UploadedFile
        if hasattr(file, 'read'):
            file_bytes = io.BytesIO(file.read())
            prs = Presentation(file_bytes)
        else:
            prs = Presentation(file)
        
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_content.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
        
        full_text = "\n\n".join(text_content)
        
        return {
            'success': True,
            'text': full_text,
            'slides': len(text_content)
        }
        
    except Exception as e:
        return {
            'success': False,
            'text': '',
            'slides': 0,
            'error': f'Error parsing PowerPoint: {str(e)}'
        }


def parse_document(file, filename: str) -> dict:
    """
    Parse a document based on its file extension.
    
    Args:
        file: File object from Streamlit uploader
        filename: Original filename to determine type
        
    Returns:
        Dictionary with parsed content or error
    """
    filename_lower = filename.lower()
    
    if filename_lower.endswith('.pdf'):
        return parse_pdf(file)
    elif filename_lower.endswith('.pptx'):
        return parse_pptx(file)
    else:
        return {
            'success': False,
            'text': '',
            'error': f'Unsupported file type. Please upload PDF or PPTX files.'
        }
